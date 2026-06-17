"""
Generates the project-continuation presentation in the same visual style as the
original Indirect Answers presentation:
  - Background:  #1B1B2F (deep dark navy)
  - Top bar:     #00B4D8 (cyan)
  - Eyebrow /
    slide tag:   #FFD66B (golden yellow, bold)
  - Body text:   #FFFFFF white / #90E0EF light-blue sub-headers
  - Muted text:  #CCCCCC
  - Font:        Calibri throughout
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import pathlib

# ── palette ──────────────────────────────────────────────────────────────────
BG          = RGBColor(0x1B, 0x1B, 0x2F)   # dark navy
CYAN        = RGBColor(0x00, 0xB4, 0xD8)   # accent bar / divider
GOLD        = RGBColor(0xFF, 0xD6, 0x6B)   # eyebrow / highlighted label
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)   # primary body text
LIGHT_BLUE  = RGBColor(0x90, 0xE0, 0xEF)  # sub-header / secondary text
GRAY        = RGBColor(0xCC, 0xCC, 0xCC)  # muted / footer
DARK_ACCENT = RGBColor(0x16, 0x3A, 0x5C)  # table header background
MID_NAVY    = RGBColor(0x22, 0x2A, 0x48)  # alternating table row tint

FONT = "Calibri"

# ── slide data ────────────────────────────────────────────────────────────────
PROJECT_TITLE  = "Beyond Indirect Answers:\nPragmatic Augmentation for CIRCA"
EYEBROW        = "Project Continuation of Louis et al. (EMNLP 2020)"
AUTHOR         = "Ron Beiden"
AUTHORS_FULL   = "Yuval Meron  ·  Ron Biden  ·  Gilad Schwartz  ·  Ziv Fenigstein"
AFFIL          = "Dept. of Industrial Engineering and Management  ·  Ben-Gurion University"
DATE           = "June 2026"
FIGURES        = pathlib.Path(__file__).parent / "figures"

MACRO_F1_ROWS = [
    ("baseline", "0.501", "0.515", "0.511"),
    ("repeat",   "0.536", "0.595", "0.597"),
    ("template", "0.545", "0.518", "0.546"),
    ("oracle",   "1.000", "1.000", "1.000"),
    ("llm",      "0.755", "0.782", "0.754"),
]

ACC_ROWS = [
    ("baseline", "0.756", "0.798", "0.804"),
    ("repeat",   "0.766", "0.814", "0.814"),
    ("template", "0.754", "0.814", "0.786"),
    ("oracle",   "1.000", "1.000", "1.000"),
    ("llm",      "0.764", "0.786", "0.760"),
]

# ── helpers ───────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


def _set_bg(slide):
    """Apply #1B1B2F solid fill to the slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _top_bar(slide, height_inches=0.12):
    """Cyan accent bar across the top of the slide."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        prs.slide_width, Inches(height_inches),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()


def _divider(slide, y_inches, width_inches=5.5):
    """Thin horizontal cyan divider, centred."""
    x = (prs.slide_width - Inches(width_inches)) / 2
    div = slide.shapes.add_shape(
        1, x, Inches(y_inches), Inches(width_inches), Inches(0.04),
    )
    div.fill.solid()
    div.fill.fore_color.rgb = CYAN
    div.line.fill.background()


def _textbox(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
    """Add a styled text box returning the shape."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name   = FONT
    f.size   = Pt(size)
    f.bold   = bold
    f.italic = italic
    f.color.rgb = color
    return txb


def _add_para(tf, text, size, color, bold=False, bullet=True, indent=0):
    """Append a paragraph to an existing text frame."""
    p   = tf.add_paragraph()
    p.level = indent
    run = p.add_run()
    run.text = ("• " if bullet else "") + text
    f = run.font
    f.name  = FONT
    f.size  = Pt(size)
    f.bold  = bold
    f.color.rgb = color
    return p


def _slide_label(slide, tag: str):
    """Small upper-right slide tag (golden)."""
    _textbox(slide, 11.0, 0.15, 2.2, 0.4, tag, 13, GOLD, bold=True, align=PP_ALIGN.RIGHT)


def _footer(slide, text: str = f"{AUTHORS_FULL}  ·  {DATE}"):
    _textbox(slide, 0.5, 7.1, 12.3, 0.35, text, 10, GRAY, align=PP_ALIGN.CENTER)


def _add_image(slide, img_path: str | pathlib.Path,
               left: float, top: float, width: float, height: float = None):
    """Insert a PNG/JPEG image. height=None → auto-scale to width."""
    img_path = pathlib.Path(img_path)
    if not img_path.exists():
        return
    pic = slide.shapes.add_picture(
        str(img_path), Inches(left), Inches(top), Inches(width),
        height=Inches(height) if height else None,
    )
    return pic


# ── slide builders ─────────────────────────────────────────────────────────────
def add_title_slide():
    slide = prs.slides.add_slide(BLANK)
    _set_bg(slide)
    _top_bar(slide, 0.14)

    # golden eyebrow
    _textbox(slide, 1.0, 1.55, 11.3, 0.7, EYEBROW, 22, GOLD, bold=True, align=PP_ALIGN.CENTER)

    # main white title
    _textbox(slide, 1.0, 2.3, 11.3, 1.8, PROJECT_TITLE, 44, WHITE, bold=True, align=PP_ALIGN.CENTER)

    _divider(slide, 4.35)

    # group members (all names)
    _textbox(slide, 1.0, 4.55, 11.3, 0.55, AUTHORS_FULL, 20, LIGHT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    _textbox(slide, 1.0, 5.1,  11.3, 0.40, AFFIL,        16, GRAY,       align=PP_ALIGN.CENTER)
    _textbox(slide, 1.0, 5.55, 11.3, 0.40, DATE,         15, GRAY,       align=PP_ALIGN.CENTER)

    _footer(slide, "")


def add_bullets_slide(title: str, bullets: list, tag: str = ""):
    slide = prs.slides.add_slide(BLANK)
    _set_bg(slide)
    _top_bar(slide)
    if tag:
        _slide_label(slide, tag)

    # golden section title
    _textbox(slide, 0.6, 0.3, 12.0, 0.75, title, 30, GOLD, bold=True)
    _divider(slide, 1.15, 12.0)

    # bullet body
    txb = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(12.0), Inches(5.6))
    tf  = txb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        if i == 0:
            p   = tf.paragraphs[0]
            run = p.add_run()
            run.text = "• " + item
            f = run.font
            f.name = FONT; f.size = Pt(24); f.color.rgb = WHITE
        else:
            _add_para(tf, item, 24, WHITE)

    _footer(slide)


def add_two_column_slide(title: str, left_title: str, left_lines: list,
                          right_title: str, right_lines: list, tag: str = ""):
    slide = prs.slides.add_slide(BLANK)
    _set_bg(slide)
    _top_bar(slide)
    if tag:
        _slide_label(slide, tag)

    _textbox(slide, 0.6, 0.3, 12.0, 0.75, title, 30, GOLD, bold=True)
    _divider(slide, 1.15, 12.0)

    # vertical separator
    sep = slide.shapes.add_shape(1, Inches(6.56), Inches(1.35), Inches(0.04), Inches(5.5))
    sep.fill.solid(); sep.fill.fore_color.rgb = CYAN; sep.line.fill.background()

    for col_x, col_title, lines in [
        (0.6,  left_title,  left_lines),
        (6.7,  right_title, right_lines),
    ]:
        _textbox(slide, col_x, 1.35, 5.7, 0.55, col_title, 22, LIGHT_BLUE, bold=True)
        txb = slide.shapes.add_textbox(Inches(col_x), Inches(2.05), Inches(5.7), Inches(4.9))
        tf  = txb.text_frame; tf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]; run = p.add_run()
                run.text = "• " + line
                f = run.font; f.name = FONT; f.size = Pt(20); f.color.rgb = WHITE
            else:
                _add_para(tf, line, 20, WHITE)

    _footer(slide)


def _cell_style(cell, text: str, size: int, color: RGBColor, bold: bool,
                bg_color: RGBColor | None = None, align=PP_ALIGN.CENTER):
    """Style a single table cell."""
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    if bg_color:
        solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
        srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", f"{bg_color[0]:02X}{bg_color[1]:02X}{bg_color[2]:02X}")

    tf  = cell.text_frame
    tf.text = text
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    f = run.font
    f.name = FONT; f.size = Pt(size); f.bold = bold; f.color.rgb = color


def add_table_slide(title: str, note: str, columns: list, rows: list, tag: str = ""):
    slide = prs.slides.add_slide(BLANK)
    _set_bg(slide)
    _top_bar(slide)
    if tag:
        _slide_label(slide, tag)

    _textbox(slide, 0.6, 0.3, 12.0, 0.75, title, 30, GOLD, bold=True)
    _divider(slide, 1.15, 12.0)
    _textbox(slide, 0.7, 1.25, 12.0, 0.45, note, 17, LIGHT_BLUE)

    n_rows = len(rows) + 1
    n_cols = len(columns)
    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.7), Inches(1.85), Inches(12.0), Inches(4.85),
    ).table

    # column widths
    col_widths = [Inches(2.0)] + [Inches(10.0 / (n_cols - 1))] * (n_cols - 1)
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = w

    # header row
    for ci, col in enumerate(columns):
        _cell_style(tbl.cell(0, ci), col, 17, WHITE, True, bg_color=DARK_ACCENT)

    # data rows
    for ri, row in enumerate(rows):
        row_bg = MID_NAVY if ri % 2 == 0 else None
        for ci, val in enumerate(row):
            bold  = ci == 0                            # mode name column
            color = GOLD if (bold and val == "llm") else (LIGHT_BLUE if bold else WHITE)
            _cell_style(tbl.cell(ri + 1, ci), val, 16, color, bold, bg_color=row_bg)

    _footer(slide)


# ── build ──────────────────────────────────────────────────────────────────────
def build_deck():
    add_title_slide()

    add_bullets_slide(
        "Why Continue This Paper?",
        [
            "Louis et al. show context helps classify indirect answers in CIRCA.",
            "Open question: can pragmatic augmentation improve robustness?",
            "We compare LLM explanation injection against controlled baselines.",
            "We expand to multiple encoders for stronger, cross-model conclusions.",
        ],
        tag="Motivation",
    )

    add_bullets_slide(
        "Research Questions",
        [
            "RQ1: Does explanation augmentation improve 6-class prediction quality?",
            "RQ2: Is LLM augmentation better than simple non-LLM controls?",
            "RQ3: Are gains consistent across BERT, MNLI-BERT, and RoBERTa?",
            "Primary metric: Macro-F1 (class-balanced).  Secondary: Accuracy.",
        ],
        tag="RQs",
    )

    add_bullets_slide(
        "Data and Task",
        [
            "Dataset: CIRCA indirect-answer corpus, 6-class pragmatic labels.",
            "Valid base examples: 30,958.  LLM-augmented subset: 5,000.",
            "Labels re-derived robustly from goldstandard1 (no leakage).",
            "Matched-cap protocol: all modes trained at the same sample budget.",
        ],
        tag="Data",
    )

    add_bullets_slide(
        "Five Augmentation Modes",
        [
            "baseline  — answer + context only (no augmentation).",
            "repeat    — duplicate input text; noise-control for extra tokens.",
            "template  — hand-written label-agnostic pragmatic hint (no label leakage).",
            "oracle    — label-description hint appended (upper-bound; leaks semantics).",
            "llm       — GPT-generated natural-language explanation per example.",
        ],
        tag="Methods",
    )

    add_bullets_slide(
        "Model Suite and Training Setup",
        [
            "Models: bert-base-uncased  ·  MNLI-BERT (textattack)  ·  roberta-base.",
            "Total runs: 15  (5 modes × 3 models), identical training budget.",
            "Framework: PyTorch + HuggingFace Transformers, 4 epochs each.",
            "All results, logs, and aggregates saved to outputs/ for traceability.",
        ],
        tag="Setup",
    )

    add_table_slide(
        "Results: Macro-F1  (Primary Metric)",
        "Higher = better.  Balanced across all 6 classes; not inflated by majority.",
        ["Mode", "BERT", "MNLI-BERT", "RoBERTa"],
        MACRO_F1_ROWS,
        tag="Results",
    )

    add_table_slide(
        "Results: Accuracy  (Secondary Metric)",
        "Accuracy can mask minority-class failures — interpret together with Macro-F1.",
        ["Mode", "BERT", "MNLI-BERT", "RoBERTa"],
        ACC_ROWS,
        tag="Results",
    )

    add_two_column_slide(
        "Key Findings",
        "What Worked",
        [
            "LLM explanations raise Macro-F1 by ~0.25 over baseline.",
            "Repeat / template controls improve modestly but far less than LLM.",
            "MNLI-BERT + LLM achieves the best Macro-F1: 0.782.",
            "Controlled design confirms pragmatic signal is the driver.",
        ],
        "What Needs Caution",
        [
            "Oracle is an artificial upper bound — not deployable.",
            "LLM gains on accuracy alone are inconsistent across models.",
            "Explanation quality variance may affect result stability.",
            "Single domain; external validity still needs study.",
        ],
        tag="Analysis",
    )

    add_bullets_slide(
        "Comparison to the Original Paper",
        [
            "Original: context window & architecture comparisons for indirect answers.",
            "This project: augmentation-focused intervention study on the same task.",
            "Added multiple controlled baselines (per professor feedback).",
            "Added cross-model validation (per professor feedback).",
            "Conclusion: pragmatic augmentation helps, especially for balanced F1.",
        ],
        tag="Discussion",
    )

    add_bullets_slide(
        "Threats to Validity and Limitations",
        [
            "Single dataset / domain limits external validity.",
            "LLM explanations may carry generation-model biases.",
            "No human scoring of explanation quality in this phase.",
            "Fixed-epoch budget — deeper tuning may shift rankings.",
        ],
        tag="Limitations",
    )

    add_bullets_slide(
        "Future Work",
        [
            "Add backtranslation and paraphrase-model augmentation as richer controls.",
            "Evaluate on additional pragmatic datasets and cross-domain transfer.",
            "Ablate explanation properties: length, structure, rationale type.",
            "Study confidence calibration and per-class error patterns.",
        ],
        tag="Future",
    )

    add_bullets_slide(
        "Reproducibility",
        [
            "Code bundle:     submission/code.zip  (src/ + paper/)",
            "Data bundle:     submission/data.zip  (circa_processed.csv + explanations)",
            "Training script: src/train.py   |   Eval orchestrator: src/evaluate.py",
            "All 15 result folders and aggregate CSVs committed under outputs/",
        ],
        tag="Repro",
    )

    # ── Figure slides ─────────────────────────────────────────────────────────
    def _fig_slide(img_file, title, note, tag):
        slide = prs.slides.add_slide(BLANK)
        _set_bg(slide)
        _top_bar(slide)
        _slide_label(slide, tag)
        _textbox(slide, 0.6, 0.3, 12.0, 0.65, title, 28, GOLD, bold=True)
        _divider(slide, 1.05, 12.0)
        _textbox(slide, 0.7, 1.12, 12.0, 0.4, note, 15, LIGHT_BLUE)
        _add_image(slide, FIGURES / img_file, left=0.9, top=1.55, width=11.4)
        _footer(slide)

    _fig_slide("macro_f1.png",
               "Macro-F1 by Augmentation Mode and Model",
               "LLM augmentation outperforms all non-oracle controls across all three encoders.",
               "Results")

    _fig_slide("accuracy.png",
               "Accuracy by Augmentation Mode and Model",
               "Accuracy gains are smaller and less consistent — Macro-F1 is the key metric.",
               "Results")

    _fig_slide("llm_gain.png",
               "LLM Augmentation Gain over Baseline  (Macro-F1)",
               "All three encoders benefit substantially from pragmatic reasoning augmentation.",
               "Analysis")

    _fig_slide("per_class_f1.png",
               "Per-Class F1: Baseline vs LLM Augmentation (BERT)",
               "Minority classes (Prob No, In the Middle) go from near-zero to >0.5 F1 with LLM.",
               "Analysis")

    # ── Closing slide ─────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(BLANK)
    _set_bg(slide)
    _top_bar(slide, 0.14)
    _textbox(slide, 1.0, 1.9, 11.3, 1.0, "Thank You", 56, WHITE, bold=True, align=PP_ALIGN.CENTER)
    _divider(slide, 3.2)
    _textbox(slide, 1.0, 3.4, 11.3, 0.6,
             "Questions & Discussion", 28, GOLD, bold=True, align=PP_ALIGN.CENTER)
    _textbox(slide, 1.0, 4.2, 11.3, 0.55, AUTHORS_FULL, 18, LIGHT_BLUE, align=PP_ALIGN.CENTER)
    _textbox(slide, 1.0, 4.75, 11.3, 0.45, AFFIL,       14, GRAY,       align=PP_ALIGN.CENTER)
    _footer(slide)


if __name__ == "__main__":
    build_deck()
    out_path = Path(__file__).resolve().parent / "Project_Continuation_Presentation.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")

